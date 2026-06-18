"""
Teacher-upload file parser (PARSE ONLY).
========================================
Teachers upload course files to S3 under teacher_uploads/<course_id>/.
This script downloads each supported file, extracts clean text into a
normalized JSON document, and uploads the JSON to parsed_json/<course_id>/.

The JSON is future input for RAG chunking — but this script does NOT chunk,
embed, load Qdrant, or touch retrieval / guardrails / AST / schemas / runtime.
It is a standalone, self-contained parser (one new file). It reuses no other
repo module (the MIT extractor has a module-level side effect, so its cleaning
heuristics are reimplemented minimally here).

Supported: .pdf .docx .txt .md .pptx .html (.htm)

Libraries (lazy-imported per format; missing → clear install message):
  PDF   -> PyMuPDF (fitz)
  DOCX  -> python-docx
  PPTX  -> python-pptx
  HTML  -> beautifulsoup4 (bs4)
  TXT/MD-> stdlib

Usage (S3):
  python data_ingestion/s3_teacher_file_parser.py \
    --bucket codingrabbit-data-dev \
    --input-prefix teacher_uploads/<course_id>/ \
    --output-prefix parsed_json/<course_id>/ \
    --profile codingrabbit-dev [--region us-east-1] [--dry-run]

Usage (local debug, no S3):
  python data_ingestion/s3_teacher_file_parser.py \
    --local-input-dir sample_teacher_uploads \
    --local-output-dir sample_parsed_json [--dry-run]
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import sys
import tempfile
from datetime import datetime, timezone
from pathlib import Path

PARSER_VERSION = "teacher_parser_v1"

SUPPORTED_EXTS = {".pdf", ".docx", ".txt", ".md", ".pptx", ".html", ".htm"}
EXT_TO_TYPE = {
    ".pdf": "pdf", ".docx": "docx", ".txt": "txt", ".md": "md",
    ".pptx": "pptx", ".html": "html", ".htm": "html",
}


# ---------------------------------------------------------------------------
# Minimal cleaning / code-detection helpers (reimplemented, not imported, to
# avoid extract_pdf_text.py's module-level os.makedirs side effect).
# ---------------------------------------------------------------------------

_CODE_KEYWORDS = {
    "int", "char", "float", "double", "void", "return", "if", "else",
    "for", "while", "do", "switch", "case", "break", "continue",
    "struct", "class", "public", "private", "virtual", "const",
    "static", "sizeof", "malloc", "free", "new", "delete",
    "printf", "scanf", "std::", "#include", "nullptr", "null",
    "strcpy", "strcat", "strlen", "def", "import", "return;",
}


def clean_text(text: str) -> str:
    """Strip control chars, collapse runs of spaces and blank lines."""
    if not text:
        return ""
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def has_code(text: str) -> bool:
    """Heuristic: fenced block, C/C++/py keyword, or brace/comment markers."""
    if not text:
        return False
    if "```" in text:
        return True
    tokens = set(re.findall(r"\b\w+\b", text.lower()))
    return bool(tokens & _CODE_KEYWORDS) or "//" in text or "{" in text


def _missing_dep(pkg: str, fmt: str) -> "SystemExit":
    return SystemExit(
        f"ERROR: {fmt} parsing requires '{pkg}', which is not installed.\n"
        f"       Install it with:  pip install {pkg}"
    )


# ---------------------------------------------------------------------------
# Block helper
# ---------------------------------------------------------------------------

def _block(block_id, block_type, text, *, page_number=None, slide_number=None,
           heading=""):
    cleaned = clean_text(text)
    return {
        "block_id": block_id,
        "block_type": block_type,
        "page_number": page_number,
        "slide_number": slide_number,
        "heading": heading,
        "text": cleaned,
        "has_code": has_code(cleaned),
    }


# ---------------------------------------------------------------------------
# Per-format parsers — each returns list[block dict]
# ---------------------------------------------------------------------------

def parse_pdf(path: Path) -> list[dict]:
    try:
        import fitz  # noqa: WPS433
    except ImportError:
        raise _missing_dep("pymupdf", "PDF")
    blocks = []
    doc = fitz.open(str(path))
    try:
        for i in range(doc.page_count):
            text = clean_text(doc[i].get_text())
            if not text:
                continue
            blocks.append(_block(f"page_{i+1}", "page", text, page_number=i + 1))
    finally:
        doc.close()
    return blocks


def parse_docx(path: Path) -> list[dict]:
    try:
        import docx  # noqa: WPS433 — python-docx
    except ImportError:
        raise _missing_dep("python-docx", "DOCX")
    document = docx.Document(str(path))

    sections: list[dict] = []
    cur_heading, cur_lines = "", []

    def flush():
        if cur_lines:
            sections.append({"heading": cur_heading, "text": "\n".join(cur_lines)})

    for para in document.paragraphs:
        txt = para.text or ""
        style = (para.style.name if para.style else "") or ""
        if style.startswith("Heading") and txt.strip():
            flush()
            cur_heading, cur_lines[:] = txt.strip(), []
        elif txt.strip():
            cur_lines.append(txt)
    flush()

    if not sections:  # no headings/paragraphs split — fall back to whole body
        whole = "\n".join(p.text for p in document.paragraphs if p.text.strip())
        if whole.strip():
            sections = [{"heading": "", "text": whole}]

    return [
        _block(f"section_{i}", "section", s["text"], heading=s["heading"])
        for i, s in enumerate(sections)
    ]


_MD_HEADING = re.compile(r"^(#{1,3})\s+(.*)$")


def parse_text(path: Path) -> list[dict]:
    """TXT/MD: split on markdown h1-h3 if present; else one block.
    Code fences are preserved verbatim (and not treated as headings)."""
    raw = path.read_text(encoding="utf-8", errors="replace")
    lines = raw.splitlines()

    # If no markdown headings, emit a single block (fences preserved as-is).
    in_fence = False
    has_heading = False
    for ln in lines:
        if ln.strip().startswith("```"):
            in_fence = not in_fence
            continue
        if not in_fence and _MD_HEADING.match(ln):
            has_heading = True
            break

    if not has_heading:
        if not raw.strip():
            return []
        return [_block("block_0", "section", raw, heading="")]

    sections: list[dict] = []
    cur_heading, cur_lines = "", []
    in_fence = False

    def flush():
        if cur_lines:
            sections.append({"heading": cur_heading, "text": "\n".join(cur_lines)})

    for ln in lines:
        if ln.strip().startswith("```"):
            in_fence = not in_fence
            cur_lines.append(ln)
            continue
        m = _MD_HEADING.match(ln) if not in_fence else None
        if m:
            flush()
            cur_heading, cur_lines[:] = m.group(2).strip(), []
        else:
            cur_lines.append(ln)
    flush()

    return [
        _block(f"section_{i}", "section", s["text"], heading=s["heading"])
        for i, s in enumerate(sections)
    ]


def parse_pptx(path: Path) -> list[dict]:
    try:
        from pptx import Presentation  # noqa: WPS433 — python-pptx
    except ImportError:
        raise _missing_dep("python-pptx", "PPTX")
    prs = Presentation(str(path))
    blocks = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        parts.append(line)
        text = clean_text("\n".join(parts))
        if not text:
            continue
        blocks.append(_block(f"slide_{idx}", "slide", text, slide_number=idx))
    return blocks


def _html_render(node) -> str:
    """Render an HTML node to text, fencing <pre> blocks. Recursive so nested
    <pre> survive (mirrors the approach used in cs50_scraper)."""
    name = getattr(node, "name", None)
    if name is None:
        return str(node).strip()
    if name == "pre":
        return f"```\n{node.get_text().rstrip(chr(10))}\n```"
    if node.find("pre"):
        return "\n".join(
            r for r in (_html_render(c) for c in node.children) if r.strip()
        )
    return node.get_text(" ", strip=True)


def parse_html(path: Path) -> list[dict]:
    try:
        from bs4 import BeautifulSoup  # noqa: WPS433
    except ImportError:
        raise _missing_dep("beautifulsoup4", "HTML")
    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="replace"),
                         "html.parser")
    for tag in soup(["nav", "script", "style", "footer", "header"]):
        tag.decompose()
    main = soup.find("main") or soup.body or soup

    headings = main.find_all(["h1", "h2", "h3"])
    if not headings:
        text = _html_render(main)
        return [_block("block_0", "section", text, heading="")] if text.strip() else []

    blocks = []
    for i, h in enumerate(headings):
        buf = []
        for sib in h.next_siblings:
            if getattr(sib, "name", None) in ("h1", "h2", "h3"):
                break
            rendered = _html_render(sib).strip()
            if rendered:
                buf.append(rendered)
        text = "\n\n".join(buf).strip()
        heading = h.get_text(" ", strip=True)
        if text or heading:
            blocks.append(_block(f"section_{i}", "section", text, heading=heading))
    return blocks


PARSERS = {
    "pdf": parse_pdf, "docx": parse_docx, "txt": parse_text, "md": parse_text,
    "pptx": parse_pptx, "html": parse_html,
}


# ---------------------------------------------------------------------------
# Envelope
# ---------------------------------------------------------------------------

def _document_id(source_key: str) -> str:
    return hashlib.sha256(source_key.encode("utf-8")).hexdigest()[:16]


def _output_stem(file_name: str, file_type: str) -> str:
    """Suffix the source type to avoid collisions (lecture1.pdf + lecture1.docx
    -> lecture1__pdf.json / lecture1__docx.json)."""
    stem = Path(file_name).stem
    return f"{stem}__{file_type}.json"


def build_envelope(*, blocks, course_id, source_uri, parsed_uri, file_name,
                   file_type, now_iso, extra_meta=None):
    meta = {
        "page_count": sum(1 for b in blocks if b["block_type"] == "page"),
        "slide_count": sum(1 for b in blocks if b["block_type"] == "slide"),
        "block_count": len(blocks),
    }
    if extra_meta:
        meta.update(extra_meta)
    return {
        "document_id": _document_id(source_uri),
        "course_id": course_id,
        "source_s3_uri": source_uri,
        "parsed_s3_uri": parsed_uri,
        "file_name": file_name,
        "file_type": file_type,
        "parser_version": PARSER_VERSION,
        "created_at": now_iso,
        "metadata": meta,
        "blocks": blocks,
    }


def _has_text(blocks) -> bool:
    return any(b["text"].strip() for b in blocks)


# ---------------------------------------------------------------------------
# Discovery
# ---------------------------------------------------------------------------

def _ext_of(name: str) -> str:
    return Path(name).suffix.lower()


def discover_s3(s3, bucket: str, prefix: str) -> list[dict]:
    """Paginated list of objects under prefix. Returns [{key, name, ext}]."""
    items = []
    token = None
    while True:
        kwargs = {"Bucket": bucket, "Prefix": prefix}
        if token:
            kwargs["ContinuationToken"] = token
        resp = s3.list_objects_v2(**kwargs)
        for obj in resp.get("Contents", []):
            key = obj["Key"]
            if key.endswith("/"):
                continue  # folder marker
            items.append({"key": key, "name": os.path.basename(key),
                          "ext": _ext_of(key)})
        if resp.get("IsTruncated"):
            token = resp.get("NextContinuationToken")
        else:
            break
    return items


def discover_local(input_dir: Path) -> list[dict]:
    items = []
    for p in sorted(input_dir.rglob("*")):
        if p.is_file():
            items.append({"key": str(p), "name": p.name, "ext": _ext_of(p.name)})
    return items


def _course_id_from_prefix(prefix: str, override: str | None) -> str:
    if override:
        return override
    parts = [p for p in prefix.strip("/").split("/") if p]
    # teacher_uploads/<course_id>/...
    if len(parts) >= 2 and parts[0] == "teacher_uploads":
        return parts[1]
    return parts[-1] if parts else "unknown_course"


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def run(args):
    now_iso = _now_iso()
    local_mode = bool(args.local_input_dir or args.local_output_dir)

    # ---- mode validation ----
    if local_mode:
        if not (args.local_input_dir and args.local_output_dir):
            raise SystemExit("ERROR: local mode needs both --local-input-dir and --local-output-dir")
        if args.bucket or args.input_prefix or args.output_prefix:
            raise SystemExit("ERROR: do not mix local mode with S3 flags (--bucket/--input-prefix/--output-prefix)")
    else:
        if not (args.bucket and args.input_prefix and args.output_prefix):
            raise SystemExit("ERROR: S3 mode needs --bucket, --input-prefix, and --output-prefix")

    # ---- discovery ----
    s3 = None
    if local_mode:
        in_dir = Path(args.local_input_dir)
        if not in_dir.is_dir():
            raise SystemExit(f"ERROR: local input dir not found: {in_dir}")
        out_dir = Path(args.local_output_dir)
        course_id = args.course_id or in_dir.name
        items = discover_local(in_dir)
    else:
        import boto3  # noqa: WPS433
        session = boto3.Session(profile_name=args.profile, region_name=args.region)
        s3 = session.client("s3")
        course_id = _course_id_from_prefix(args.input_prefix, args.course_id)
        items = discover_s3(s3, args.bucket, args.input_prefix)

    print(f"Teacher file parser ({PARSER_VERSION}){' — DRY RUN' if args.dry_run else ''}")
    print(f"  mode: {'local' if local_mode else 'S3'} | course_id: {course_id}")
    print(f"  discovered {len(items)} object(s)\n")

    supported = [it for it in items if it["ext"] in SUPPORTED_EXTS]
    skipped = [it for it in items if it["ext"] not in SUPPORTED_EXTS]

    parsed_count = 0
    failed = []
    skipped_empty = []
    total_blocks = 0
    output_uris = []

    tmp_ctx = tempfile.TemporaryDirectory()
    tmp_dir = Path(tmp_ctx.name)

    try:
        for it in supported:
            file_type = EXT_TO_TYPE[it["ext"]]
            out_name = _output_stem(it["name"], file_type)
            if local_mode:
                source_uri = f"file://{Path(it['key']).resolve()}"
                parsed_uri = f"file://{(Path(args.local_output_dir).resolve() / out_name)}"
            else:
                source_uri = f"s3://{args.bucket}/{it['key']}"
                parsed_uri = f"s3://{args.bucket}/{args.output_prefix.rstrip('/')}/{out_name}"

            if args.dry_run:
                print(f"  [{file_type:<4}] {it['name']}  ->  {parsed_uri}")
                continue

            try:
                # obtain a local path to the source file
                if local_mode:
                    local_path = Path(it["key"])
                else:
                    local_path = tmp_dir / it["name"]
                    s3.download_file(args.bucket, it["key"], str(local_path))

                blocks = PARSERS[file_type](local_path)

                if not _has_text(blocks):
                    print(f"  ❌ {it['name']}: no text extracted (not uploaded)")
                    skipped_empty.append(it["name"])
                    failed.append(it["name"])
                    continue

                envelope = build_envelope(
                    blocks=blocks, course_id=course_id, source_uri=source_uri,
                    parsed_uri=parsed_uri, file_name=it["name"],
                    file_type=file_type, now_iso=now_iso,
                )

                # write JSON to temp, then place at destination
                tmp_json = tmp_dir / out_name
                tmp_json.write_text(json.dumps(envelope, ensure_ascii=False, indent=2),
                                    encoding="utf-8")
                if local_mode:
                    out_dir.mkdir(parents=True, exist_ok=True)
                    (out_dir / out_name).write_text(
                        tmp_json.read_text(encoding="utf-8"), encoding="utf-8")
                else:
                    s3.upload_file(
                        str(tmp_json), args.bucket,
                        f"{args.output_prefix.rstrip('/')}/{out_name}",
                        ExtraArgs={"ServerSideEncryption": "AES256",
                                   "ContentType": "application/json"},
                    )

                parsed_count += 1
                total_blocks += len(blocks)
                output_uris.append(parsed_uri)
                print(f"  ✓ [{file_type}] {it['name']} -> {len(blocks)} blocks -> {out_name}")

            except SystemExit:
                raise  # missing-dependency message: surface immediately
            except Exception as e:  # noqa: BLE001 — fail loud per file, keep going
                print(f"  ❌ {it['name']} FAILED: {e}")
                failed.append(it["name"])
    finally:
        tmp_ctx.cleanup()

    # ---- summary ----
    print("\n=== Summary ===")
    print(f"  discovered: {len(items)}")
    print(f"  parsed:     {parsed_count}")
    print(f"  skipped (unsupported): {len(skipped)}" +
          (f" -> {[s['name'] for s in skipped]}" if skipped else ""))
    print(f"  failed:     {len(failed)}" + (f" -> {failed}" if failed else ""))
    if skipped_empty:
        print(f"  empty (no text, not uploaded): {skipped_empty}")
    print(f"  total blocks extracted: {total_blocks}")
    if output_uris:
        print("  output:")
        for u in output_uris:
            print(f"    {u}")
    if args.dry_run:
        print("\nDry run complete. Nothing downloaded or uploaded.")


def main():
    p = argparse.ArgumentParser(description="Parse teacher-uploaded course files into normalized JSON (parse only).")
    # S3 mode
    p.add_argument("--bucket")
    p.add_argument("--input-prefix")
    p.add_argument("--output-prefix")
    p.add_argument("--profile", default=None, help="AWS profile (S3 mode)")
    p.add_argument("--region", default="us-east-1", help="AWS region (S3 mode)")
    # local mode
    p.add_argument("--local-input-dir")
    p.add_argument("--local-output-dir")
    # shared
    p.add_argument("--course-id", default=None, help="Override inferred course_id")
    p.add_argument("--dry-run", action="store_true",
                   help="List files + planned output paths; download/upload nothing")
    run(p.parse_args())


if __name__ == "__main__":
    main()
